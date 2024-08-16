import os
import io
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

import os.path
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError

SCOPES = [
    'https://www.googleapis.com/auth/drive.readonly',
    'https://www.googleapis.com/auth/drive.file',
    'https://www.googleapis.com/auth/drive'
]
# Google Drive API authentication
def authenticate_drive_api():
  creds = None
  # The file token.json stores the user's access and refresh tokens, and is
  # created automatically when the authorization flow completes for the first
  # time.
  if os.path.exists("token.json"):
    creds = Credentials.from_authorized_user_file("token.json", SCOPES)
  # If there are no (valid) credentials available, let the user log in.
  if not creds or not creds.valid:
    if creds and creds.expired and creds.refresh_token:
      creds.refresh(Request())
    else:
      flow = InstalledAppFlow.from_client_secrets_file(
          "/Users/dingshengliu/Desktop/ChatbotAI/Testing/ReadFiles/cred.json", SCOPES
      )
      creds = flow.run_local_server(port=0)
    # Save the credentials for the next run
    with open("token.json", "w") as token:
      token.write(creds.to_json())

  try:
    service = build("drive", "v3", credentials=creds)

    # Call the Drive v3 API
    results = (
        service.files()
        .list(pageSize=10, fields="nextPageToken, files(id, name)")
        .execute()
    )
    items = results.get("files", [])

    if not items:
      print("No files found.")
      return
    print("Files:")
    for item in items:
      print(f"{item['name']} ({item['id']})")
    return service
  except HttpError as error:
    # TODO(developer) - Handle errors from drive API.
    print(f"An error occurred: {error}")

# Function to download file from Google Drive
def download_file(service, file_id, file_name):
    try:
        request = service.files().get_media(fileId=file_id)
        fh = io.FileIO(file_name, 'wb')
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while done is False:
            status, done = downloader.next_chunk()
        fh.close()
    except Exception as e:
        raise Exception(f"Error downloading file {file_name}: {e}")

# Function to export Google Docs files to a downloadable format
def export_google_doc(service, file_id, mime_type, file_name):
    try:
        request = service.files().export_media(fileId=file_id, mimeType=mime_type)
        fh = io.FileIO(file_name, 'wb')
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while done is False:
            status, done = downloader.next_chunk()
        fh.close()
    except Exception as e:
        raise Exception(f"Error exporting file {file_name}: {e}")

# Function to extract text from PDF
def extract_text_from_pdf(file_path):
    text = ""
    try:
        reader = PdfReader(file_path)
        for page in reader.pages:
            text += page.extract_text()
    except Exception as e:
        raise Exception(f"Error reading PDF file {file_path}: {e}")
    return text

# Function to extract text from DOCX
def extract_text_from_docx(file_path):
    try:
        doc = docx.Document(file_path)
        text = [paragraph.text for paragraph in doc.paragraphs]
    except Exception as e:
        raise Exception(f"Error reading DOCX file {file_path}: {e}")
    return '\n'.join(text)

# Function to extract text from PPTX
def extract_text_from_pptx(file_path):
    text = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
    except Exception as e:
        raise Exception(f"Error reading PPTX file {file_path}: {e}")
    return '\n'.join(text)

# Function to resolve the original file ID if the item is a shortcut
def resolve_shortcut(service, file_id):
    try:
        file = service.files().get(fileId=file_id, fields='id, mimeType, shortcutDetails').execute()
        if file.get('mimeType') == 'application/vnd.google-apps.shortcut':
            return file['shortcutDetails']['targetId']
        return file_id
    except Exception as e:
        raise Exception(f"Error resolving shortcut for file {file_id}: {e}")

# Main function to process all files in a Google Drive folder
def process_drive_folder(service, folder_id, output_file):
    query = f"'{folder_id}' in parents"
    results = service.files().list(q=query, pageSize=1000).execute()
    items = results.get('files', [])

    with open(output_file, 'w', encoding='utf-8') as output:
        for item in items:
            file_id = item['id']
            file_name = item['name']
            print(f"Processing file: {file_name}")

            file_id = resolve_shortcut(service, file_id)
            mime_type = item.get('mimeType')
            text = ""
            try:
                if mime_type == 'application/vnd.google-apps.document':
                    export_google_doc(service, file_id, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', file_name + '.docx')
                    text = extract_text_from_docx(file_name + '.docx')
                elif mime_type == 'application/vnd.google-apps.spreadsheet':
                    export_google_doc(service, file_id, 'application/pdf', file_name + '.pdf')
                    text = extract_text_from_pdf(file_name + '.pdf')
                elif mime_type == 'application/vnd.google-apps.presentation':
                    export_google_doc(service, file_id, 'application/vnd.openxmlformats-officedocument.presentationml.presentation', file_name + '.pptx')
                    text = extract_text_from_pptx(file_name + '.pptx')
                else:
                    download_file(service, file_id, file_name)
                    file_extension = os.path.splitext(file_name)[1].lower()

                    if file_extension == '.pdf':
                        text = extract_text_from_pdf(file_name)
                    elif file_extension == '.docx':
                        text = extract_text_from_docx(file_name)
                    elif file_extension == '.pptx':
                        text = extract_text_from_pptx(file_name)
                    else:
                        text = f"Unsupported file type: {file_name}\n"
            except Exception as e:
                text = f"Error processing file {file_name}: {e}\n"
                print(text)
                
            output.write(text + '\n\n')
            
            # Clean up temporary files
            for ext in ['.pdf', '.docx', '.pptx']:
                temp_file = file_name + ext
                if os.path.exists(temp_file):
                    os.remove(temp_file)
            if os.path.exists(file_name):
                os.remove(file_name)
                
if __name__ == '__main__':
    service = authenticate_drive_api()
    folder_id = '1t6qqKO9Yy1bHbefcX38MpEI2AyE0zlOq'
    output_file = '/Users/dingshengliu/Desktop/ChatbotAI/Testing/data/output.txt'
    process_drive_folder(service, folder_id, output_file)